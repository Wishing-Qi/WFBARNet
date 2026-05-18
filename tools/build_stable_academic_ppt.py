# -*- coding: utf-8 -*-
from __future__ import annotations

import sys
from pathlib import Path
from typing import Iterable

from PIL import Image


ROOT = Path(__file__).resolve().parents[1]
LOCAL_DEPS = ROOT / ".codex_deps" / "python_pptx"
if LOCAL_DEPS.exists():
    sys.path.insert(0, str(LOCAL_DEPS))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = ROOT / "outputs"
OUT_PATH = OUT_DIR / "WFBARNet_答辩PPT_稳重学术型.pptx"
LOGO_SRC = ROOT / "tools" / "hylogo.png"
LOGO_TRANSPARENT = OUT_DIR / "ppt_assets" / "hylogo_transparent.png"


SLIDE_W = 13.333333
SLIDE_H = 7.5

NAVY = "132B5E"
BLUE = "2E7BC2"
CYAN = "46C8E9"
PURPLE = "8A55E6"
BG = "F6F8FC"
BG2 = "EAF2FA"
PANEL = "FFFFFF"
LINE = "C8D6EA"
TEXT = "152238"
MUTED = "64748B"
PALE = "EDF4FA"

FONT_TITLE = "Noto Serif SC"
FONT_BODY = "Microsoft YaHei"


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_text(run, size: int, color: str = TEXT, bold: bool = False, font: str = FONT_BODY) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 20,
    color: str = TEXT,
    bold: bool = False,
    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_text(run, size, color, bold, font)
    return shape


def add_multiline(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: Iterable[str],
    size: int = 18,
    color: str = TEXT,
    bullet: bool = True,
    line_spacing: float = 1.15,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ("• " if bullet else "") + line
        p.font.name = FONT_BODY
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(6 * line_spacing)
    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = PANEL,
    line: str = LINE,
    radius: bool = True,
    transparency: int = 0,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.fill.transparency = transparency
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(0.8)
    return shp


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: str = LINE, width: float = 1.2):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    return conn


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: str = CYAN, width: float = 2.0):
    line = add_line(slide, x1, y1, x2, y2, color, width)
    # Some python-pptx builds expose arrowhead enums differently. A plain connector is safer;
    # the terminal marker is added as a small triangle for consistent rendering.
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x2 - 0.08), Inches(y2 - 0.08), Inches(0.16), Inches(0.16))
    tri.fill.solid()
    tri.fill.fore_color.rgb = rgb(color)
    tri.line.color.rgb = rgb(color)
    if abs(x2 - x1) >= abs(y2 - y1):
        tri.rotation = 90 if x2 >= x1 else 270
    else:
        tri.rotation = 180 if y2 >= y1 else 0
    return line


def make_logo_transparent() -> Path:
    LOGO_TRANSPARENT.parent.mkdir(parents=True, exist_ok=True)
    img = Image.open(LOGO_SRC).convert("RGBA")
    px = img.load()
    w, h = img.size
    for x in range(w):
        for y in range(h):
            r, g, b, a = px[x, y]
            if a < 8 or (r > 238 and g > 238 and b > 238):
                px[x, y] = (r, g, b, 0)
    img.save(LOGO_TRANSPARENT)
    return LOGO_TRANSPARENT


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(BG)


def add_theme(slide, page_no: int, footer: str = "WFBARNet 本科毕业论文答辩") -> None:
    set_background(slide)
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.88), Inches(SLIDE_H))
    left.fill.solid()
    left.fill.fore_color.rgb = rgb(NAVY)
    left.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.88), Inches(0), Inches(0.06), Inches(SLIDE_H))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(CYAN)
    accent.line.fill.background()

    add_line(slide, 1.45, 0.74, 12.42, 0.74, LINE, 1.0)
    add_line(slide, 1.45, 6.80, 12.42, 6.80, LINE, 0.8)
    add_textbox(slide, 1.44, 7.02, 3.3, 0.2, footer, 8, MUTED)
    add_textbox(slide, 12.07, 7.02, 0.25, 0.2, f"{page_no:02d}", 8, MUTED, align=PP_ALIGN.RIGHT)

    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(11.20), Inches(0), Inches(2.14), Inches(1.48))
    tri.fill.solid()
    tri.fill.fore_color.rgb = rgb(PURPLE)
    tri.fill.transparency = 82
    tri.line.fill.background()
    tri.rotation = 180


def add_logo(slide, x: float = 11.86, y: float = 0.25, size: float = 0.62) -> None:
    slide.shapes.add_picture(str(LOGO_TRANSPARENT), Inches(x), Inches(y), width=Inches(size), height=Inches(size))


def add_title(slide, title: str, subtitle: str, page_no: int) -> None:
    add_theme(slide, page_no)
    add_logo(slide)
    add_textbox(slide, 1.44, 0.38, 6.8, 0.34, title, 22, TEXT, True, FONT_TITLE)
    add_textbox(slide, 1.44, 0.78, 7.6, 0.24, subtitle, 9, MUTED)
    bar1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.44), Inches(1.08), Inches(1.38), Inches(0.055))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = rgb(CYAN)
    bar1.line.fill.background()
    bar2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.92), Inches(1.08), Inches(0.72), Inches(0.055))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = rgb(PURPLE)
    bar2.line.fill.background()


def info_chip(slide, x: float, y: float, text: str, color: str = BLUE, w: float | None = None):
    w = w or max(0.7, 0.14 * len(text) + 0.22)
    shp = add_rect(slide, x, y, w, 0.34, fill="F4FAFE", line="D5E6F4")
    add_textbox(slide, x + 0.08, y + 0.08, w - 0.16, 0.13, text, 7, color, True, align=PP_ALIGN.CENTER)
    return shp


def add_card_title(slide, x: float, y: float, title: str, num: str | None = None):
    if num is not None:
        box = add_rect(slide, x, y, 0.44, 0.32, NAVY, NAVY)
        add_textbox(slide, x, y + 0.07, 0.44, 0.1, num, 8, "FFFFFF", True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.58, y, 3.6, 0.25, title, 15, TEXT, True)
    else:
        add_textbox(slide, x, y, 3.8, 0.25, title, 15, TEXT, True)


def build_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_theme(slide, 1, "稳重学术型")
    slide.shapes.add_picture(str(LOGO_TRANSPARENT), Inches(10.35), Inches(2.58), width=Inches(1.62), height=Inches(1.62))
    add_textbox(slide, 1.52, 1.34, 3.5, 0.32, "本科毕业论文答辩", 13, BLUE)
    add_textbox(slide, 1.52, 1.94, 5.6, 0.58, "WFBARNet", 36, NAVY, False, FONT_TITLE)
    add_textbox(slide, 1.52, 2.74, 8.7, 0.42, "面向羽毛球视频的本地智能分析系统", 21, TEXT, True, FONT_TITLE)
    add_textbox(slide, 1.52, 3.20, 3.5, 0.38, "设计与实现", 20, TEXT, False, FONT_TITLE)
    add_line(slide, 1.52, 3.78, 4.62, 3.78, CYAN, 2.8)
    add_line(slide, 4.70, 3.78, 5.62, 3.78, PURPLE, 2.8)
    rows = [("答辩人", "沃寅博"), ("专业", "软件工程"), ("学院", "江苏师范大学科文学院"), ("答辩时间", "5月11日")]
    y = 4.46
    for k, v in rows:
        add_textbox(slide, 1.55, y, 0.8, 0.20, f"{k}：", 10, MUTED)
        add_textbox(slide, 2.45, y, 3.2, 0.20, v, 11, TEXT, True)
        y += 0.42
    add_textbox(slide, 1.52, 6.28, 7.0, 0.22, "轨迹跟踪 · 姿态估计 · 球场映射 · 回合统计 · 可视化复盘", 9, MUTED)


def build_agenda(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "目录", "论文主要章节结构", 2)
    items = [
        ("01", "绪论", "研究背景、意义与国内外现状"),
        ("02", "系统需求与总体设计", "功能需求、总体架构与模块划分"),
        ("03", "关键算法与实现", "轨迹跟踪、姿态估计、球场映射、事件识别"),
        ("04", "实验设计与结果分析", "测试数据、指标对比与可视化结果"),
        ("05", "总结与展望", "研究结论、创新点与后续优化方向"),
    ]
    y = 1.63
    for num, title, desc in items:
        add_rect(slide, 1.68, y, 9.8, 0.66, "FFFFFF", "DDE8F4")
        add_rect(slide, 1.93, y + 0.17, 0.55, 0.32, NAVY, NAVY)
        add_textbox(slide, 1.93, y + 0.245, 0.55, 0.10, num, 8, "FFFFFF", True, align=PP_ALIGN.CENTER)
        add_textbox(slide, 2.72, y + 0.11, 2.6, 0.20, title, 13, TEXT, True)
        add_textbox(slide, 2.72, y + 0.39, 4.2, 0.14, desc, 8, MUTED)
        add_line(slide, 9.78, y + 0.34, 10.80, y + 0.34, CYAN, 1.8)
        y += 0.86


def build_background_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "项目背景与问题", "羽毛球训练复盘需要可量化、可回放、可解释的数据支撑", 3)
    add_textbox(slide, 1.50, 1.45, 4.7, 0.3, "传统复盘痛点", 18, NAVY, True, FONT_TITLE)
    pain = [
        "人工观察依赖经验，难以稳定记录球路、站位与移动负荷。",
        "羽毛球目标小、速度快、模糊和遮挡频繁，检测难度高。",
        "单独检测球或人不足以支持训练分析，需要球、球员、球场联合理解。",
        "训练过程需要保留可追溯日志，方便排查误检、漏检和模型问题。",
    ]
    add_multiline(slide, 1.55, 1.95, 5.1, 3.4, pain, 14)
    add_rect(slide, 7.05, 1.45, 4.9, 4.45, "FFFFFF", LINE)
    add_textbox(slide, 7.38, 1.75, 3.8, 0.25, "WFBARNet 的应对思路", 17, NAVY, True, FONT_TITLE)
    concepts = [("球", "TrackNetV3 轨迹跟踪"), ("球员", "YOLO Pose 姿态估计"), ("球场", "场线检测与单应性映射")]
    x = 7.42
    for label, desc in concepts:
        add_rect(slide, x, 2.45, 1.22, 1.20, "F4FAFE", "D5E6F4")
        add_textbox(slide, x, 2.72, 1.22, 0.25, label, 22, BLUE, True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x - 0.05, 3.82, 1.32, 0.42, desc, 8, MUTED, align=PP_ALIGN.CENTER)
        x += 1.45
    add_arrow(slide, 8.00, 4.70, 10.85, 4.70, CYAN, 2.2)
    add_textbox(slide, 7.46, 5.05, 4.15, 0.35, "目标：将视频画面转化为结构化训练数据", 15, TEXT, True, align=PP_ALIGN.CENTER)


def build_position_goal(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "项目定位与总体目标", "本地化、可解释、可调试、可扩展的辅助分析系统", 4)
    cards = [
        ("输入", ["本地视频", "摄像头实时流", "批量视频帧"], BLUE),
        ("核心处理", ["球轨迹", "人体姿态", "球场映射", "轨迹事件"], NAVY),
        ("输出", ["可视化叠加", "回合统计", "JSON/CSV/NPY", "调试日志"], CYAN),
        ("应用", ["训练复盘", "技战术观察", "数据整理", "模型实验"], PURPLE),
    ]
    x = 1.42
    for title, lines, color in cards:
        add_rect(slide, x, 1.68, 2.55, 3.85, "FFFFFF", LINE)
        add_rect(slide, x + 0.22, 1.95, 0.58, 0.42, color, color)
        add_textbox(slide, x + 0.96, 2.00, 1.3, 0.24, title, 16, TEXT, True, FONT_TITLE)
        add_multiline(slide, x + 0.30, 2.72, 2.00, 1.95, lines, 13, TEXT)
        x += 2.85
    add_rect(slide, 1.65, 5.95, 10.0, 0.52, "EEF6FB", "D5E6F4")
    add_textbox(slide, 1.90, 6.10, 9.45, 0.18, "定位说明：系统用于辅助分析和训练复盘，不作为自动裁判或唯一判定依据。", 12, NAVY, True)


def build_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "系统总体架构", "PyQt6 应用层、视觉推理层、后处理统计层和数据导出层协同工作", 5)
    layers = [
        ("应用层", "PyQt6 桌面界面 · 视频播放 · 摄像头实时分析 · 数据页展示"),
        ("视觉推理层", "TrackBranch · PoseBranch · CourtLineDetector · BSTStrokeRecognizer"),
        ("后处理层", "轨迹滤波 · 姿态稳定 · 事件检测 · 球员距离累计 · 回合统计"),
        ("输出层", "视频叠加 · 标准球场视图 · JSON/CSV/NPY · JSONL/CSV 调试日志"),
    ]
    y = 1.45
    for idx, (name, desc) in enumerate(layers):
        add_rect(slide, 1.52, y, 10.55, 0.78, "FFFFFF", LINE)
        add_rect(slide, 1.80, y + 0.18, 1.05, 0.42, [NAVY, BLUE, CYAN, PURPLE][idx], [NAVY, BLUE, CYAN, PURPLE][idx])
        add_textbox(slide, 1.80, y + 0.285, 1.05, 0.1, name, 10, "FFFFFF", True, align=PP_ALIGN.CENTER)
        add_textbox(slide, 3.08, y + 0.20, 7.9, 0.24, desc, 13, TEXT, True if idx == 1 else False)
        if idx < len(layers) - 1:
            add_arrow(slide, 6.66, y + 0.78, 6.66, y + 1.05, CYAN, 1.8)
        y += 1.13
    add_textbox(slide, 1.54, 6.30, 10.6, 0.22, "结构特点：核心算法模块与界面展示解耦，便于离线 runner、实时分析和后续模型替换。", 11, MUTED)


def build_core_pipeline(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "核心流程：从视频帧到分析结果", "每一帧经过多路推理、统一封装，再进入事件和统计模块", 6)
    steps = [
        ("视频帧", "读取帧与时间戳"),
        ("三帧窗口", "[prev, current, next]"),
        ("多路推理", "轨迹 / 姿态 / 球场"),
        ("FrameResult", "pose + track"),
        ("事件识别", "hit / landing / out_of_frame"),
        ("统计展示", "UI payload + 日志"),
    ]
    x = 1.35
    y = 2.05
    for idx, (title, desc) in enumerate(steps):
        add_rect(slide, x, y, 1.55, 1.05, "FFFFFF", LINE)
        add_textbox(slide, x + 0.12, y + 0.18, 1.30, 0.22, title, 13, NAVY, True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.10, y + 0.55, 1.35, 0.22, desc, 8, MUTED, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_arrow(slide, x + 1.55, y + 0.52, x + 1.95, y + 0.52, CYAN, 1.8)
        x += 1.95
    add_rect(slide, 1.55, 4.15, 10.3, 1.45, "EEF6FB", "D5E6F4")
    add_textbox(slide, 1.90, 4.43, 9.6, 0.25, "运行时主循环", 16, TEXT, True, FONT_TITLE)
    details = [
        "TrackNet 使用三帧窗口捕捉短时运动信息；姿态和球场模块可按不同频率运行。",
        "RealtimeTrajectoryEventDetector 不依赖 UI 渲染帧率，保证事件日志连续。",
        "统一输出结构为 FrameResult(frame_id, pose, track)，便于后续 UI、日志和统计复用。",
    ]
    add_multiline(slide, 1.92, 4.86, 9.3, 0.56, details, 10)


def build_track_module(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "羽毛球轨迹跟踪模块", "解决小目标、高速运动、短时丢失和多候选误检问题", 7)
    add_rect(slide, 1.45, 1.50, 5.2, 4.65, "FFFFFF", LINE)
    add_card_title(slide, 1.75, 1.80, "TrackNetV3 风格轨迹分支")
    add_multiline(
        slide,
        1.82,
        2.35,
        4.35,
        2.2,
        [
            "输入连续三帧，拼接为 9 通道时序窗口。",
            "模型输出羽毛球热力图，后处理解码候选球点。",
            "默认输入尺寸 512 × 288，候选点最多保留 5 个。",
            "支持 PyTorch 权重与 TensorRT engine 后端。",
        ],
        12,
    )
    add_textbox(slide, 1.82, 5.30, 4.4, 0.22, "预处理：resize → BGR/RGB → 9 通道张量 → 热力图还原", 10, MUTED)
    add_rect(slide, 7.05, 1.50, 4.85, 4.65, "FFFFFF", LINE)
    add_card_title(slide, 7.35, 1.80, "候选点解码流程")
    labels = ["原始帧", "热力图", "连通域", "最终球点"]
    x = 7.48
    for idx, label in enumerate(labels):
        add_rect(slide, x, 2.55, 0.94, 0.78, "F4FAFE", "D5E6F4")
        add_textbox(slide, x, 2.83, 0.94, 0.12, label, 8, NAVY, True, align=PP_ALIGN.CENTER)
        if idx < 3:
            add_arrow(slide, x + 0.94, 2.94, x + 1.20, 2.94, CYAN, 1.5)
        x += 1.34
    add_multiline(slide, 7.42, 3.92, 3.9, 1.15, ["连通域峰值、均值、面积和紧致度综合排序。", "通过缩放比例将热力图坐标还原到原始帧。"], 11)


def build_filter_module(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "轨迹滤波与连续性修复", "把模型候选点转换为稳定、可解释的逐帧轨迹", 8)
    add_rect(slide, 1.50, 1.42, 5.25, 4.75, "FFFFFF", LINE)
    add_textbox(slide, 1.82, 1.78, 4.2, 0.24, "TrackNetV3TrajectoryFilter", 16, NAVY, True, FONT_TITLE)
    add_multiline(
        slide,
        1.85,
        2.35,
        4.45,
        1.8,
        [
            "可见高置信候选点优先保留。",
            "过滤低置信、越界和明显不可靠候选。",
            "fixed-lag 模式可对中间缺失段做线性 inpaint。",
            "保留可插拔接口，便于替换滤波策略。",
        ],
        12,
    )
    for i, field in enumerate(["action", "reason", "output_x/y", "inpaint_mask", "relock"]):
        info_chip(slide, 1.85 + (i % 3) * 1.42, 4.74 + (i // 3) * 0.44, field, BLUE, 1.15)
    add_rect(slide, 7.15, 1.42, 4.72, 4.75, "FFFFFF", LINE)
    add_textbox(slide, 7.48, 1.78, 3.3, 0.24, "轨迹修复示意", 16, NAVY, True, FONT_TITLE)
    chart_x, chart_y = 7.55, 2.45
    add_line(slide, chart_x, chart_y + 2.35, chart_x + 3.70, chart_y + 2.35, LINE, 1.0)
    add_line(slide, chart_x, chart_y, chart_x, chart_y + 2.35, LINE, 1.0)
    pts = [(chart_x + 0.2, chart_y + 1.95), (chart_x + 0.85, chart_y + 1.35), (chart_x + 1.45, chart_y + 0.80), (chart_x + 2.15, chart_y + 1.12), (chart_x + 2.86, chart_y + 0.55), (chart_x + 3.45, chart_y + 0.92)]
    for a, b in zip(pts, pts[1:]):
        add_line(slide, a[0], a[1], b[0], b[1], CYAN, 2.0)
    for idx, (px, py) in enumerate(pts):
        color = PURPLE if idx == 3 else NAVY
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px - 0.055), Inches(py - 0.055), Inches(0.11), Inches(0.11))
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(color)
        dot.line.color.rgb = rgb(color)
    add_textbox(slide, 7.58, 5.23, 3.8, 0.22, "紫色点表示缺失段修复或重新锁定后的轨迹点。", 10, MUTED)


def build_pose_module(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "球员姿态估计与稳定", "从球的位置扩展到球员行为、站位与运动负荷分析", 9)
    add_rect(slide, 1.50, 1.52, 5.10, 4.55, "FFFFFF", LINE)
    add_card_title(slide, 1.80, 1.84, "姿态估计输出")
    add_multiline(
        slide,
        1.84,
        2.38,
        4.28,
        2.0,
        [
            "主后端为 Ultralytics YOLO Pose，同时保留 MMPose 接入能力。",
            "输出人体框、关键点、关键点置信度和球员身份索引。",
            "用于绘制骨架、估计脚下锚点、区分上方/下方球员。",
        ],
        12,
    )
    # Simple skeleton symbol
    cx, cy = 4.88, 4.55
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.15), Inches(cy - 0.82), Inches(0.30), Inches(0.30))
    head.fill.solid(); head.fill.fore_color.rgb = rgb(CYAN); head.line.color.rgb = rgb(CYAN)
    for x1, y1, x2, y2 in [
        (cx, cy - 0.50, cx, cy + 0.10), (cx, cy - 0.28, cx - 0.42, cy - 0.05),
        (cx, cy - 0.28, cx + 0.42, cy - 0.05), (cx, cy + 0.10, cx - 0.35, cy + 0.65),
        (cx, cy + 0.10, cx + 0.35, cy + 0.65),
    ]:
        add_line(slide, x1, y1, x2, y2, BLUE, 2.0)
    add_rect(slide, 7.05, 1.52, 4.90, 4.55, "FFFFFF", LINE)
    add_card_title(slide, 7.36, 1.84, "稳定与投影")
    for idx, (title, desc) in enumerate([("关键点", "人体骨架"), ("脚下锚点", "bbox / 脚踝"), ("场地坐标", "image_to_court_h")]):
        y = 2.45 + idx * 1.02
        add_rect(slide, 7.46, y, 1.35, 0.55, [NAVY, BLUE, CYAN][idx], [NAVY, BLUE, CYAN][idx])
        add_textbox(slide, 7.46, y + 0.18, 1.35, 0.10, title, 9, "FFFFFF", True, align=PP_ALIGN.CENTER)
        add_textbox(slide, 9.08, y + 0.11, 2.15, 0.20, desc, 12, TEXT, True)
        if idx < 2:
            add_arrow(slide, 8.13, y + 0.55, 8.13, y + 0.86, CYAN, 1.6)
    add_textbox(slide, 7.48, 5.36, 4.0, 0.2, "结合球场位置和历史结果，降低抖动、漏检和身份交换影响。", 10, MUTED)


def build_court_module(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "球场线检测与场地坐标映射", "单应性矩阵把图像像素转换为真实场地空间", 10)
    add_rect(slide, 1.48, 1.50, 4.85, 4.65, "FFFFFF", LINE)
    add_card_title(slide, 1.78, 1.82, "统一检测接口")
    add_multiline(
        slide,
        1.82,
        2.35,
        3.95,
        1.65,
        [
            "默认后端：OpenCV 球场线检测。",
            "可选后端：MonoTrack 风格传统 CV。",
            "输出角点、标准模板投影和 image_to_court_h。",
        ],
        12,
    )
    info_chip(slide, 1.82, 4.60, "create_court_line_detector", BLUE, 2.05)
    info_chip(slide, 4.02, 4.60, "predict_court_lines", PURPLE, 1.62)
    add_rect(slide, 7.00, 1.50, 4.95, 4.65, "FFFFFF", LINE)
    add_card_title(slide, 7.30, 1.82, "坐标映射示意")
    # Image plane quadrilateral
    x0, y0 = 7.52, 2.55
    pts = [(x0, y0 + 1.4), (x0 + 1.75, y0 + 0.75), (x0 + 1.35, y0), (x0 - 0.32, y0 + 0.45)]
    for a, b in zip(pts, pts[1:] + pts[:1]):
        add_line(slide, a[0], a[1], b[0], b[1], BLUE, 2.0)
    add_textbox(slide, x0 - 0.20, y0 + 1.72, 1.7, 0.16, "图像平面", 10, MUTED, align=PP_ALIGN.CENTER)
    add_arrow(slide, 9.45, 3.23, 10.00, 3.23, CYAN, 2.2)
    # Standard court
    court = add_rect(slide, 10.20, 2.18, 1.25, 2.35, "F4FAFE", BLUE, False)
    add_line(slide, 10.20, 3.35, 11.45, 3.35, BLUE, 1.1)
    add_line(slide, 10.62, 2.18, 10.62, 4.53, BLUE, 0.8)
    add_line(slide, 11.03, 2.18, 11.03, 4.53, BLUE, 0.8)
    add_textbox(slide, 10.05, 4.82, 1.65, 0.16, "标准场地坐标", 10, MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.52, 5.35, 3.95, 0.22, "用途：球场线叠加、热力图、击球区域、移动距离统计。", 10, MUTED)


def build_event_module(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "轨迹事件识别", "规则式识别击球、落点和出画事件，便于解释与调试", 11)
    add_rect(slide, 1.48, 1.48, 5.65, 4.62, "FFFFFF", LINE)
    add_card_title(slide, 1.78, 1.80, "事件轨迹示意")
    chart_x, chart_y = 1.85, 2.45
    add_line(slide, chart_x, chart_y + 2.35, chart_x + 4.55, chart_y + 2.35, LINE, 1)
    add_line(slide, chart_x, chart_y, chart_x, chart_y + 2.35, LINE, 1)
    pts = [(chart_x + 0.15, chart_y + 1.9), (chart_x + 0.85, chart_y + 1.15), (chart_x + 1.55, chart_y + 0.55), (chart_x + 2.35, chart_y + 1.08), (chart_x + 3.05, chart_y + 1.72), (chart_x + 3.75, chart_y + 2.03)]
    for a, b in zip(pts, pts[1:]):
        add_line(slide, a[0], a[1], b[0], b[1], CYAN, 2.0)
    events = [(pts[2], "hit", "DC2626"), (pts[5], "landing", "16A34A"), ((chart_x + 4.25, chart_y + 0.35), "out", PURPLE)]
    for (px, py), label, color in events:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px - 0.07), Inches(py - 0.07), Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = rgb(color); dot.line.color.rgb = rgb(color)
        add_textbox(slide, px - 0.24, py + 0.14, 0.55, 0.15, label, 8, color, True, align=PP_ALIGN.CENTER)
    add_rect(slide, 7.45, 1.48, 4.45, 4.62, "FFFFFF", LINE)
    add_card_title(slide, 7.76, 1.80, "规则与过滤")
    rules = [
        ("hit", "vy_reversal / vx_reversal", "主击球规则"),
        ("landing", "speed_step / speed_drop", "落点候选"),
        ("out_of_frame", "visibility_drop_edge", "出画/跟踪丢失"),
    ]
    y = 2.45
    for event, rule, note in rules:
        add_rect(slide, 7.78, y, 3.45, 0.62, "F8FBFE", "DDE8F4")
        add_textbox(slide, 7.95, y + 0.10, 0.95, 0.16, event, 10, NAVY, True)
        add_textbox(slide, 8.95, y + 0.10, 1.55, 0.16, rule, 9, BLUE, True)
        add_textbox(slide, 8.95, y + 0.34, 1.78, 0.14, note, 8, MUTED)
        y += 0.88
    add_textbox(slide, 7.80, 5.32, 3.8, 0.22, "加速度峰值和速度峰值只作为辅助证据，降低误检。", 10, MUTED)


def build_bst_module(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "击球动作识别与 BST 输入", "在击球事件附近截取时序片段，进一步识别动作类别", 12)
    steps = [("hit_event", "轨迹事件触发"), ("时序片段", "姿态 + 球点 + 场地位置"), ("BST 模型", "动作分类"), ("输出", "类别 / 置信度 / Top-K")]
    x = 1.62
    for idx, (title, desc) in enumerate(steps):
        add_rect(slide, x, 2.10, 2.05, 1.18, "FFFFFF", LINE)
        add_textbox(slide, x + 0.18, 2.38, 1.7, 0.24, title, 14, NAVY, True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.18, 2.78, 1.7, 0.16, desc, 9, MUTED, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_arrow(slide, x + 2.05, 2.69, x + 2.50, 2.69, CYAN, 1.8)
        x += 2.55
    add_rect(slide, 1.58, 4.30, 4.95, 1.45, "FFFFFF", LINE)
    add_textbox(slide, 1.92, 4.58, 4.2, 0.22, "可选能力", 15, TEXT, True, FONT_TITLE)
    add_multiline(slide, 1.92, 4.96, 4.12, 0.52, ["没有 BST 权重时，系统仍可完成轨迹、姿态、球场和回合统计。", "BSTInputBuilder 可导出后续训练或实验所需数据。"], 10)
    add_rect(slide, 7.05, 4.30, 4.85, 1.45, "EEF6FB", "D5E6F4")
    add_textbox(slide, 7.38, 4.58, 4.0, 0.22, "输出示例", 15, TEXT, True, FONT_TITLE)
    for i, chip in enumerate(["杀球 0.92", "挑球 0.76", "放网 0.69"]):
        info_chip(slide, 7.40 + i * 1.32, 5.06, chip, [NAVY, BLUE, PURPLE][i], 1.05)


def build_rally_stats(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "回合统计与数据指标", "把逐帧视觉结果汇总为可用于训练复盘的回合级指标", 13)
    add_rect(slide, 1.46, 1.42, 4.92, 4.80, "FFFFFF", LINE)
    add_card_title(slide, 1.78, 1.76, "回合状态")
    timeline = [("未开始", NAVY), ("回合中", BLUE), ("回合结束", CYAN)]
    y = 2.42
    for idx, (label, color) in enumerate(timeline):
        add_rect(slide, 1.88, y, 1.15, 0.48, color, color)
        add_textbox(slide, 1.88, y + 0.15, 1.15, 0.10, label, 9, "FFFFFF", True, align=PP_ALIGN.CENTER)
        if idx < 2:
            add_arrow(slide, 3.08, y + 0.24, 3.68, y + 0.24, CYAN, 1.6)
        y += 0.86
    add_textbox(slide, 1.88, 5.15, 3.85, 0.30, "开始：稳定球飞行段或第一条 hit；结束：可信 landing。", 10, MUTED)
    add_rect(slide, 6.82, 1.42, 5.10, 4.80, "FFFFFF", LINE)
    add_card_title(slide, 7.14, 1.76, "核心指标")
    metrics = [
        ("回合", "时长 / 击球次数 / 落点 / 出画"),
        ("运动", "距离 / 平均速度 / 最大速度"),
        ("动作", "启动 / 急停 / 高强度移动"),
        ("区域", "前场 / 中场 / 后场击球分布"),
        ("质量", "球点 / 姿态 / 球场有效率"),
    ]
    y = 2.30
    for head, desc in metrics:
        add_rect(slide, 7.20, y, 0.72, 0.36, NAVY, NAVY)
        add_textbox(slide, 7.20, y + 0.11, 0.72, 0.08, head, 8, "FFFFFF", True, align=PP_ALIGN.CENTER)
        add_textbox(slide, 8.12, y + 0.08, 3.05, 0.14, desc, 11, TEXT)
        y += 0.62


def build_ui_export(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "可视化界面与数据导出", "将算法结果转化为可查看、可复盘、可调试的工作台", 14)
    add_rect(slide, 1.42, 1.44, 6.10, 4.72, "FFFFFF", LINE)
    add_card_title(slide, 1.72, 1.76, "PyQt6 可视化界面")
    # UI mockup
    add_rect(slide, 1.75, 2.20, 3.65, 2.35, "0F172A", NAVY, False)
    add_textbox(slide, 3.58, 3.28, 0.8, 0.16, "视频区", 12, "FFFFFF", True, align=PP_ALIGN.CENTER)
    add_rect(slide, 5.62, 2.20, 1.35, 2.35, "F4FAFE", "D5E6F4")
    add_textbox(slide, 5.72, 2.38, 1.12, 0.16, "数据面板", 9, NAVY, True, align=PP_ALIGN.CENTER)
    for i in range(4):
        add_rect(slide, 5.78, 2.78 + i * 0.38, 1.00, 0.20, "FFFFFF", "DDE8F4", radius=False)
    add_multiline(slide, 1.82, 4.92, 4.8, 0.55, ["球点、轨迹尾迹、人体骨架、球场线、事件 marker 同步展示。", "支持播放控制、日志开关、标准球场视图和热力图。"], 10)
    add_rect(slide, 7.90, 1.44, 4.08, 4.72, "FFFFFF", LINE)
    add_card_title(slide, 8.18, 1.76, "数据导出与调试")
    for idx, item in enumerate(["JSON", "CSV", "NPY", "JSONL", "Debug CSV"]):
        info_chip(slide, 8.20 + (idx % 2) * 1.60, 2.36 + (idx // 2) * 0.52, item, [NAVY, BLUE, CYAN, PURPLE, BLUE][idx], 1.25)
    add_textbox(slide, 8.22, 4.10, 3.2, 0.20, "典型日志", 13, TEXT, True)
    add_multiline(slide, 8.22, 4.45, 3.25, 0.70, ["*_frame_log.jsonl：逐帧球点、姿态、事件。", "*_track_debug.csv：候选选择、滤波动作、inpaint、relock。"], 10)


def build_tech_stack(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "技术栈与工程实现", "Python 桌面应用架构，兼顾深度学习推理、传统 CV 与可维护性", 15)
    stack = [
        ("应用层", "PyQt6 · 视频播放 · 实时交互"),
        ("模型层", "PyTorch · TrackNetV3 · YOLO Pose · BST"),
        ("视觉层", "OpenCV · NumPy · 球场线检测 · 单应性映射"),
        ("导出层", "JSON · CSV · JSONL · NPY · 可视化视频"),
        ("测试层", "轨迹事件 · 回合统计 · 球场检测 · 导出器"),
    ]
    y = 1.42
    for idx, (head, body) in enumerate(stack):
        add_rect(slide, 1.50, y, 5.05, 0.62, "FFFFFF", LINE)
        add_rect(slide, 1.76, y + 0.15, 0.72, 0.32, [NAVY, BLUE, CYAN, PURPLE, NAVY][idx], [NAVY, BLUE, CYAN, PURPLE, NAVY][idx])
        add_textbox(slide, 1.76, y + 0.23, 0.72, 0.08, head, 8, "FFFFFF", True, align=PP_ALIGN.CENTER)
        add_textbox(slide, 2.70, y + 0.18, 3.45, 0.15, body, 11, TEXT)
        y += 0.82
    add_rect(slide, 7.05, 1.42, 4.85, 4.70, "FFFFFF", LINE)
    add_card_title(slide, 7.36, 1.76, "关键目录")
    dirs = ["src/models", "src/postprocess", "src/court", "src/runners", "apps/pyqt6", "tests"]
    for idx, item in enumerate(dirs):
        info_chip(slide, 7.45 + (idx % 2) * 1.65, 2.38 + (idx // 2) * 0.58, item, BLUE if idx % 2 == 0 else PURPLE, 1.35)
    add_textbox(slide, 7.46, 5.18, 3.75, 0.28, "工程特点：模块边界清晰，离线 runner 与实时 UI 共用核心组件。", 11, MUTED)


def build_summary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "总结与展望", "研究结论、创新点与后续优化方向", 16)
    add_rect(slide, 1.50, 1.48, 5.05, 3.85, "FFFFFF", LINE)
    add_card_title(slide, 1.82, 1.82, "研究结论")
    conclusions = [
        "构建了围绕羽毛球、球员、球场三对象的本地分析链路。",
        "实现了轨迹事件识别与回合级数据统计，支持训练复盘。",
        "提供 PyQt6 可视化界面与调试日志，便于算法验证和迭代。",
    ]
    add_multiline(slide, 1.86, 2.40, 4.22, 1.7, conclusions, 12)
    add_rect(slide, 7.00, 1.48, 4.90, 3.85, "FFFFFF", LINE)
    add_card_title(slide, 7.32, 1.82, "创新点与展望")
    future = [
        "可解释事件规则、多源数据融合、可追溯调试体系。",
        "引入多回合自动切分，增强击球归属和界内外判断。",
        "提升 TensorRT 实时性能和跨场景泛化能力。",
    ]
    add_multiline(slide, 7.36, 2.40, 4.05, 1.7, future, 12)
    add_textbox(slide, 6.65, 5.94, 1.8, 0.36, "谢谢聆听", 23, NAVY, False, FONT_TITLE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 5.52, 6.38, 4.05, 0.18, "欢迎各位老师批评指正", 11, MUTED, align=PP_ALIGN.CENTER)


def build_ppt() -> Path:
    make_logo_transparent()
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    builders = [
        build_cover,
        build_agenda,
        build_background_problem,
        build_position_goal,
        build_architecture,
        build_core_pipeline,
        build_track_module,
        build_filter_module,
        build_pose_module,
        build_court_module,
        build_event_module,
        build_bst_module,
        build_rally_stats,
        build_ui_export,
        build_tech_stack,
        build_summary,
    ]
    for builder in builders:
        builder(prs)

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    print(build_ppt())
